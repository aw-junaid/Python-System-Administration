#!/usr/bin/env python3
"""
create_powerpoint.py
======================
Build a PowerPoint slide deck automatically with python-pptx from a
simple CSV outline (one row per slide: Title, Body).

Usage
-----
    python create_powerpoint.py --file outline.csv --output deck.pptx

If --file is omitted, a sample outline CSV is generated first.

Expected output
----------------
A .pptx file with:
    - A title slide (deck title + subtitle)
    - One "Title and Content" slide per row of the outline CSV, with
      the row's "Body" text split into bullet points on "|"

Requirements
------------
    pip install python-pptx
"""

import argparse
import csv
import os

from pptx import Presentation
from pptx.util import Inches, Pt


def create_sample_outline(path: str) -> None:
    content = (
        "Title,Body\n"
        "Introduction,Welcome to the quarterly review|Goals for this presentation\n"
        "Key Results,Revenue grew 18% QoQ|Customer churn down 4%|New markets opened: 2\n"
        "Next Steps,Expand into APAC|Hire 5 new engineers|Launch v2.0 in Q3\n"
    )
    with open(path, "w", encoding="utf-8") as f:
        f.write(content)
    print(f"[info] No --file given, created a sample outline CSV at: {path}")


def build_presentation(outline_csv: str, output: str, deck_title: str) -> None:
    if not os.path.exists(outline_csv):
        print(f"[error] File not found: {outline_csv}")
        return

    prs = Presentation()

    # --- Title slide ---
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = deck_title
    slide.placeholders[1].text = "Generated automatically with python-pptx"

    # --- Content slides from CSV outline ---
    content_layout = prs.slide_layouts[1]  # "Title and Content"
    slide_count = 0

    with open(outline_csv, newline="", encoding="utf-8") as f:
        reader = csv.DictReader(f)
        for row in reader:
            slide = prs.slides.add_slide(content_layout)
            slide.shapes.title.text = row.get("Title", f"Slide {slide_count + 1}")

            body_placeholder = slide.placeholders[1]
            tf = body_placeholder.text_frame
            tf.clear()

            bullets = [b.strip() for b in row.get("Body", "").split("|") if b.strip()]
            for i, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(20)

            slide_count += 1

    prs.save(output)
    print(f"[success] Presentation written to: {output}")
    print(f"          Slides created: {slide_count + 1} (1 title slide + {slide_count} content slides)")


def main():
    parser = argparse.ArgumentParser(description="Build a PowerPoint deck automatically from a CSV outline.")
    parser.add_argument("--file", default=None, help="Path to the outline CSV (columns: Title, Body)")
    parser.add_argument("--output", default="deck.pptx", help="Path for the output .pptx file")
    parser.add_argument("--title", default="Quarterly Review", help="Title for the deck's title slide")
    args = parser.parse_args()

    path = args.file
    if path is None:
        path = "outline.csv"
        create_sample_outline(path)

    build_presentation(path, args.output, args.title)


if __name__ == "__main__":
    main()
